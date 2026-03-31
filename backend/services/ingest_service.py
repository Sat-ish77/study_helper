"""
backend/services/ingest_service.py
File ingestion — extracts text, chunks, embeds, stores in Supabase.
Ported from study_helper/ingest.py with Streamlit removed.

Supports: PDF (with Mistral OCR), DOCX, PPTX
Embeddings: text-embedding-3-small → 1536 dims (NEVER CHANGE)
Storage: sh_document_chunks table
"""

import os
import io
from datetime import datetime, timezone
from openai import OpenAI
from supabase import create_client
from dotenv import load_dotenv

load_dotenv()

EMBEDDING_MODEL = "text-embedding-3-small"
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 150


def _get_supabase():
    return create_client(
        os.getenv("SUPABASE_URL"),
        os.getenv("SUPABASE_SERVICE_ROLE_KEY")
    )


def _get_openai():
    return OpenAI(api_key=os.getenv("OPENAI_API_KEY"))


def embed_texts(texts: list[str]) -> list[list[float]]:
    """Embed a batch of texts. Returns list of 1536-dim vectors."""
    client = _get_openai()
    response = client.embeddings.create(
        model=EMBEDDING_MODEL,
        input=texts
    )
    return [item.embedding for item in response.data]


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks."""
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk_words = words[i:i + chunk_size]
        chunks.append(" ".join(chunk_words))
        i += chunk_size - overlap
    return [c for c in chunks if len(c.strip()) > 50]


async def process_file(
    file_bytes: bytes,
    filename: str,
    content_type: str,
    user_id: str,
) -> dict:
    """
    Main ingestion pipeline:
    1. Extract text from file
    2. Chunk the text
    3. Embed each chunk
    4. Store in sh_document_chunks
    """
    # Delete existing chunks for this file+user (re-upload)
    sb = _get_supabase()
    sb.table("sh_document_chunks")\
        .delete()\
        .eq("user_id", user_id)\
        .eq("source_file", filename)\
        .execute()

    # Extract text
    text_chunks = await extract_text(file_bytes, filename, content_type)

    if not text_chunks:
        return {"filename": filename, "chunks_created": 0, "status": "no text extracted"}

    # Embed in batches of 20
    batch_size = 20
    rows = []

    for i in range(0, len(text_chunks), batch_size):
        batch = text_chunks[i:i + batch_size]
        texts = [c["content"] for c in batch]
        embeddings = embed_texts(texts)

        for j, (chunk, embedding) in enumerate(zip(batch, embeddings)):
            rows.append({
                "user_id": user_id,
                "content": chunk["content"],
                "source_file": filename,
                "filetype": content_type.split("/")[-1],
                "page_num": chunk.get("page_num"),
                "slide_num": chunk.get("slide_num"),
                "embedding": embedding,
            })

    # Insert all chunks
    if rows:
        sb.table("sh_document_chunks").insert(rows).execute()

    return {
        "filename": filename,
        "chunks_created": len(rows),
        "status": "success"
    }


async def extract_text(
    file_bytes: bytes,
    filename: str,
    content_type: str
) -> list[dict]:
    """Extract text from PDF, DOCX, or PPTX."""

    if "pdf" in content_type:
        return await _extract_pdf(file_bytes, filename)
    elif "wordprocessingml" in content_type:
        return _extract_docx(file_bytes, filename)
    elif "presentationml" in content_type:
        return _extract_pptx(file_bytes, filename)
    else:
        return []


async def _extract_pdf(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract PDF text using Mistral OCR (primary) or pypdf (fallback)."""
    chunks = []

    # Try Mistral OCR first
    mistral_key = os.getenv("MISTRAL_API_KEY")
    if mistral_key:
        try:
            import base64
            import httpx

            b64 = base64.b64encode(file_bytes).decode()
            async with httpx.AsyncClient() as client:
                resp = await client.post(
                    "https://api.mistral.ai/v1/ocr",
                    headers={"Authorization": f"Bearer {mistral_key}"},
                    json={
                        "model": "mistral-ocr-latest",
                        "document": {
                            "type": "document_url",
                            "document_url": f"data:application/pdf;base64,{b64}"
                        }
                    },
                    timeout=60
                )
                if resp.status_code == 200:
                    data = resp.json()
                    for page in data.get("pages", []):
                        page_text = page.get("markdown", "")
                        page_num = page.get("index", 0)
                        for chunk in chunk_text(page_text):
                            chunks.append({"content": chunk, "page_num": page_num + 1})
                    if chunks:
                        return chunks
        except Exception as e:
            print(f"[ingest] Mistral OCR failed, falling back: {e}")

    # Fallback to pypdf
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            for chunk in chunk_text(text):
                chunks.append({"content": chunk, "page_num": page_num})
    except Exception as e:
        print(f"[ingest] pypdf failed: {e}")

    return chunks


def _extract_docx(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return [{"content": chunk, "page_num": None} for chunk in chunk_text(full_text)]
    except Exception as e:
        print(f"[ingest] DOCX extraction failed: {e}")
        return []


def _extract_pptx(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract text from PPTX slide by slide."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        chunks = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slide_text = "\n".join(texts)
            for chunk in chunk_text(slide_text):
                chunks.append({"content": chunk, "slide_num": slide_num})
        return chunks
    except Exception as e:
        print(f"[ingest] PPTX extraction failed: {e}")
        return []