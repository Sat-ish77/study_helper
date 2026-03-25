"""
ingest.py — Study Helper v2
Per-user document ingestion → Supabase pgvector.
Supports: PDF, DOCX, PPTX, TXT, PNG, JPG, WEBP (via Mistral OCR)
"""
from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import List

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from dotenv import load_dotenv

load_dotenv()

CHUNK_SIZE    = 1000
CHUNK_OVERLAP = 150

SUPPORTED_TYPES = {"pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg", "webp"}


# ── Loaders ───────────────────────────────────────────────────────────────────

def load_pdf(path: Path) -> List[Document]:
    from langchain_community.document_loaders import PyPDFLoader
    loader = PyPDFLoader(str(path))
    docs   = loader.load()
    for d in docs:
        d.metadata["filename"] = path.name
        d.metadata["filetype"] = "pdf"
    return docs


def load_docx(path: Path) -> List[Document]:
    from langchain_community.document_loaders import Docx2txtLoader
    loader = Docx2txtLoader(str(path))
    docs   = loader.load()
    for d in docs:
        d.metadata["filename"] = path.name
        d.metadata["filetype"] = "docx"
    return docs


def load_pptx(path: Path) -> List[Document]:
    from pptx import Presentation
    prs  = Presentation(str(path))
    docs = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        text = "\n".join(parts).strip()
        if text:
            docs.append(Document(
                page_content=text,
                metadata={"filename": path.name, "filetype": "pptx", "slide": i}
            ))
    return docs


def load_txt(path: Path) -> List[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [Document(
        page_content=text,
        metadata={"filename": path.name, "filetype": "txt"}
    )]


def load_image(path: Path) -> List[Document]:
    """
    Extract text from image using Mistral OCR API.
    No system install needed — pure API call.
    Falls back to pytesseract if MISTRAL_API_KEY not set.
    """
    mistral_key = os.getenv("MISTRAL_API_KEY")

    # ── Mistral OCR (preferred) ───────────────────────────────────────────────
    if mistral_key:
        try:
            import base64
            from mistralai import Mistral

            client = Mistral(api_key=mistral_key)

            with open(path, "rb") as f:
                image_data = base64.b64encode(f.read()).decode()

            ext  = path.suffix.lower().lstrip(".")
            mime = {
                "jpg":  "image/jpeg",
                "jpeg": "image/jpeg",
                "png":  "image/png",
                "webp": "image/webp",
            }.get(ext, "image/png")

            response = client.ocr.process(
                model="mistral-ocr-latest",
                document={
                    "type":      "image_url",
                    "image_url": f"data:{mime};base64,{image_data}",
                }
            )

            text = "\n".join(
                page.markdown for page in response.pages
            ).strip()

            if not text:
                return []

            return [Document(
                page_content=text,
                metadata={"filename": path.name, "filetype": "image"}
            )]

        except Exception as e:
            print(f"[Mistral OCR] Failed on {path.name}: {e}")
            return []

    # ── pytesseract fallback (if no Mistral key) ──────────────────────────────
    try:
        import pytesseract
        from PIL import Image
        img  = Image.open(str(path))
        text = pytesseract.image_to_string(img).strip()
        if not text:
            return []
        return [Document(
            page_content=text,
            metadata={"filename": path.name, "filetype": "image"}
        )]
    except ImportError:
        return []
    except Exception as e:
        print(f"[pytesseract OCR] Failed on {path.name}: {e}")
        return []


def load_file(path: Path) -> List[Document]:
    ext = path.suffix.lower().lstrip(".")
    if ext == "pdf":
        return load_pdf(path)
    elif ext == "docx":
        return load_docx(path)
    elif ext == "pptx":
        return load_pptx(path)
    elif ext == "txt":
        return load_txt(path)
    elif ext in {"png", "jpg", "jpeg", "webp"}:
        return load_image(path)
    return []


# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_documents(docs: List[Document]) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
    )
    return splitter.split_documents(docs)


# ── Embeddings ────────────────────────────────────────────────────────────────

def _get_embeddings():
    """Always use OpenAI embeddings — reliable, cheap (~$0.0002/PDF)."""
    from langchain_openai import OpenAIEmbeddings
    return OpenAIEmbeddings(model="text-embedding-3-small")


# ── Supabase pgvector insert ──────────────────────────────────────────────────

def insert_chunks(chunks: List[Document], user_id: str) -> int:
    from supabase import create_client

    supabase   = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_ANON_KEY"))
    embeddings = _get_embeddings()

    texts  = [c.page_content for c in chunks]
    embeds = embeddings.embed_documents(texts)

    rows = []
    for chunk, embed in zip(chunks, embeds):
        meta = chunk.metadata
        rows.append({
            "user_id":     user_id,
            "content":     chunk.page_content,
            "embedding":   embed,
            "source_file": meta.get("filename", ""),
            "filetype":    meta.get("filetype", ""),
            "page_num":    meta.get("page"),
            "slide_num":   meta.get("slide"),
        })

    batch_size = 50
    for i in range(0, len(rows), batch_size):
        supabase.table("sh_document_chunks").insert(rows[i:i+batch_size]).execute()

    return len(rows)


def delete_user_document(user_id: str, filename: str) -> bool:
    from supabase import create_client
    try:
        supabase = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_ANON_KEY"))
        supabase.table("sh_document_chunks") \
            .delete() \
            .eq("user_id", user_id) \
            .eq("source_file", filename) \
            .execute()
        return True
    except Exception as e:
        print(f"[delete_user_document] {e}")
        return False


def get_user_documents(user_id: str) -> List[str]:
    from supabase import create_client
    try:
        supabase = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_ANON_KEY"))
        res = supabase.table("sh_document_chunks") \
            .select("source_file") \
            .eq("user_id", user_id) \
            .execute()
        files = set(r["source_file"] for r in (res.data or []) if r.get("source_file"))
        return sorted(files)
    except Exception:
        return []


# ── Main entry ────────────────────────────────────────────────────────────────

def ingest_uploaded_file(uploaded_file, user_id: str) -> dict:
    ext = uploaded_file.name.rsplit(".", 1)[-1].lower()
    if ext not in SUPPORTED_TYPES:
        return {"success": False, "error": f".{ext} is not supported."}

    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
        tmp.write(uploaded_file.getbuffer())
        tmp_path = Path(tmp.name)

    try:
        docs = load_file(tmp_path)
        for d in docs:
            d.metadata["filename"] = uploaded_file.name

        if not docs:
            return {"success": False, "error": "No text could be extracted from this file."}

        chunks = chunk_documents(docs)
        count  = insert_chunks(chunks, user_id)

        return {
            "success":  True,
            "filename": uploaded_file.name,
            "pages":    len(docs),
            "chunks":   count,
        }
    except Exception as e:
        return {"success": False, "error": str(e)}
    finally:
        tmp_path.unlink(missing_ok=True)


# ── Streamlit upload widget ───────────────────────────────────────────────────

def render_upload_widget(user_id: str):
    import streamlit as st

    uploaded_files = st.file_uploader(
        "Drop files here — PDF, DOCX, PPTX, TXT, or images",
        type=["pdf", "docx", "pptx", "txt", "png", "jpg", "jpeg", "webp"],
        accept_multiple_files=True,
        key=f"uploader_{user_id}",
    )

    if uploaded_files:
        for uf in uploaded_files:
            with st.spinner(f"Processing {uf.name}..."):
                result = ingest_uploaded_file(uf, user_id)
            if result["success"]:
                st.success(
                    f"✓ **{result['filename']}** — "
                    f"{result['pages']} pages, {result['chunks']} chunks indexed"
                )
            else:
                st.error(f"✗ {uf.name}: {result['error']}")

    docs = get_user_documents(user_id)
    if docs:
        with st.expander(f"Your documents ({len(docs)})", expanded=True):
            for doc in docs:
                col1, col2 = st.columns([5, 1])
                col1.markdown(f"📄 {doc}")
                if col2.button("Remove", key=f"del_{doc}_{user_id}"):
                    if delete_user_document(user_id, doc):
                        st.success(f"Removed {doc}")
                        st.rerun()
    else:
        st.info("No documents uploaded yet. Upload your study materials above.")